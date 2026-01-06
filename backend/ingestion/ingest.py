# backend/ingestion/ingest.py
import time
import os
import shutil
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from datetime import datetime, timezone
from processing.preprocess import clean_text, chunk_text, save_chunks
from ingestion.store import add_metadata
import hashlib
from pptx import Presentation
from pypdf import PdfReader


def extract_text_from_pptx(path):
    try:
        prs = Presentation(path)
        texts = []

        for idx, slide in enumerate(prs.slides):
            slide_text = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text.append(shape.text)

            if slide_text:
                texts.append("\n".join(slide_text))

            # mark slide boundary
            texts.append(f"\n--- SLIDE {idx+1} BREAK ---\n")

        return "\n".join(texts)

    except Exception as e:
        print("⚠ PPTX parse failed:", e)
        return ""

    try:
        prs = Presentation(path)
        texts = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)

        return "\n".join(texts)

    except Exception as e:
        print("⚠ PPTX parse failed:", e)
        return ""

def extract_text_from_pdf(path):
    try:
        reader = PdfReader(path)
        text = []
        for page in reader.pages:
            text.append(page.extract_text() or "")
        return "\n".join(text)
    except Exception as e:
        return f"(pdf extraction error: {e})"

def compute_file_hash(path):
    h = hashlib.sha256()
    with open(path, "rb") as f:
        while chunk := f.read(8192):
            h.update(chunk)
    return h.hexdigest()


WATCH_DIR = "ingest_queue"
PROCESSED_DIR = "processed"

os.makedirs(WATCH_DIR, exist_ok=True)
os.makedirs(PROCESSED_DIR, exist_ok=True)


def process_file(path):
    file_hash = compute_file_hash(path)

    meta = {
        "filename": os.path.basename(path),
        "hash": file_hash,
        "path": os.path.abspath(path),
        "size_bytes": os.path.getsize(path),
        "ingested_at": datetime.now(timezone.utc).isoformat()
    }
    print(">>> PROCESSING:", meta)

    # add metadata to store
    add_metadata(meta)

    lower = path.lower()
    if lower.endswith(".txt"):
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            print("Text content (first 500 chars):\n", text[:500])
        except Exception as e:
            print("Error reading txt:", e)

    elif lower.endswith(".pdf"):
        text = extract_text_from_pdf(path)
        print("PDF text (first 1000 chars):\n", text[:1000])

    elif lower.endswith(".pptx"):
        text = extract_text_from_pptx(path)
        print("PPTX text (first 1000 chars):\n", text[:1000])

    else:
        # For images/audio/others we currently just store metadata.
        print("Non-text file detected. Metadata stored. (extend for OCR later)")

    if not text or not text.strip():
        print("⚠ No extractable text in this file, skipping chunking")
        return

    cleaned = clean_text(text)
    chunks = chunk_text(cleaned)
    print("CHUNKS CREATED:", len(chunks))

    save_chunks(meta["filename"], meta["hash"], chunks)

    # move processed file
    dest = os.path.join(PROCESSED_DIR, os.path.basename(path))
    try:
        shutil.move(path, dest)
        print(f"Moved to processed: {dest}")
    except Exception as e:
        print("Error moving file:", e)

class IngestHandler(FileSystemEventHandler):
    def on_created(self, event):
        # ignore directories
        if event.is_directory:
            return
        # small delay to ensure file write completes
        time.sleep(0.2)
        print(f"[{datetime.now(timezone.utc).isoformat()}] New file detected:", event.src_path)
        process_file(event.src_path)

if __name__ == "__main__":
    print("Starting watcher on:", os.path.abspath(WATCH_DIR))
    observer = Observer()
    observer.schedule(IngestHandler(), WATCH_DIR, recursive=False)
    observer.start()
    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        observer.stop()
    observer.join()
